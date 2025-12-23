import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Theme Colors
BG_COLOR = RGBColor(10, 10, 10)       # #0a0a0a
TEXT_PRIMARY = RGBColor(255, 255, 255) # #ffffff
TEXT_SECONDARY = RGBColor(160, 160, 160) # #a0a0a0
ACCENT_GOLD = RGBColor(212, 175, 55)   # #D4AF37
ACCENT_HOVER = RGBColor(241, 196, 15)  # #F1C40F

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def style_text_run(run, font_size, color=TEXT_PRIMARY, bold=False):
    run.font.size = Pt(font_size)
    run.font.name = 'Arial' # Close to Inter
    run.font.color.rgb = color
    run.font.bold = bold

def create_button(slide, text, left, top, width, height, primary=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    fill = shape.fill
    fill.solid()
    
    line = shape.line
    
    if primary:
        fill.fore_color.rgb = ACCENT_GOLD
        line.color.rgb = ACCENT_GOLD
        text_color = RGBColor(0, 0, 0)
    else:
        fill.background() # Transparent-ish
        line.color.rgb = ACCENT_GOLD
        line.width = Pt(2)
        text_color = ACCENT_GOLD # Gold text

    tf = shape.text_frame
    tf.text = text
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = 'Arial'

def create_presentation():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "Batsman Pro"
    p.alignment = PP_ALIGN.CENTER
    style_text_run(p.runs[0], 60, ACCENT_GOLD, True)

    # Subtitle
    top = Inches(4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Unlock Your Potential With AI-Powered Cricket Analytics"
    p.alignment = PP_ALIGN.CENTER
    style_text_run(p.runs[0], 24, TEXT_PRIMARY)

    # --- Slide 2: About / Hero ---
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Main Heading
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(4.5)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Unlock Your Potential"
    style_text_run(p.runs[0], 40, TEXT_PRIMARY, True)
    
    # "With Batsman Pro"
    top = Inches(1.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "With Batsman Pro"
    style_text_run(p.runs[0], 40, ACCENT_GOLD, True)

    # Description
    top = Inches(2.5)
    width = Inches(5)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "The ultimate tool for cricketers. Detect bat-ball contact, generate automatic highlights, and classify shots with professional-grade AI."
    style_text_run(p.runs[0], 18, TEXT_SECONDARY)

    # Buttons
    create_button(slide, "Download APK", Inches(0.5), Inches(4.5), Inches(2), Inches(0.6), primary=True)
    create_button(slide, "Get on Play Store", Inches(2.8), Inches(4.5), Inches(2.2), Inches(0.6), primary=False)

    # Placeholder Image (Mockup)
    # Using a rectangle to represent the phone
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(1), Inches(3), Inches(5.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(30, 30, 30)
    shape.line.color.rgb = RGBColor(60, 60, 60)
    
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.text = "[App Screen]"
    p.alignment = PP_ALIGN.CENTER
    style_text_run(p.runs[0], 14, TEXT_SECONDARY)


    # --- Slide 3: Key Features ---
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Heading
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Key Features"
    p.alignment = PP_ALIGN.CENTER
    style_text_run(p.runs[0], 36, ACCENT_GOLD, True)

    # Grid of Features
    features = [
        ("AI Detection", "State-of-the-art bat-ball contact detection using YOLO based models."),
        ("Highlight Generation", "Automatically condenses hours of footage into the key moments."),
        ("Shot Classification", "Identifies whether it's a cover drive, pull shot, or defense."),
        ("Secure Processing", "Your data is processed securely with enterprise-grade privacy standards.")
    ]

    start_x = 0.5
    start_y = 1.8
    w = 4.2
    h = 2
    margin_x = 0.5
    margin_y = 0.5

    for i, (f_title, f_desc) in enumerate(features):
        row = i // 2
        col = i % 2
        x = Inches(start_x + col * (w + margin_x))
        y = Inches(start_y + row * (h + margin_y))

        # Card Background
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(20, 20, 20) # Slightly lighter than bg
        shape.line.fill.background() # No border

        # Title
        txBox = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.1), Inches(w-0.2), Inches(0.5))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = f_title
        style_text_run(p.runs[0], 18, ACCENT_GOLD, True)

        # Desc
        txBox = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.6), Inches(w-0.2), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = f_desc
        style_text_run(p.runs[0], 14, TEXT_SECONDARY)


    # --- Slide 4: How It Works ---
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Heading
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "How It Works"
    p.alignment = PP_ALIGN.CENTER
    style_text_run(p.runs[0], 36, ACCENT_GOLD, True)

    steps = [
        ("1", "Upload Video", "Upload your match footage securely."),
        ("2", "AI Detection", "Our advanced AI detects bat-ball contact."),
        ("3", "Highlight Gen", "Watch auto-generated highlights instantly."),
        ("4", "Shot Analytics", "Get deep insights into every shot played.")
    ]
    
    # Horizontal Layout
    card_w = 2.2
    gap = 0.2
    start_x = (10 - (4 * card_w + 3 * gap)) / 2 # Center align
    
    for i, (num, title, desc) in enumerate(steps):
        x = Inches(start_x + i * (card_w + gap))
        y = Inches(2.5)
        
        # Step Number Circle
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.8), y - Inches(0.8), Inches(0.6), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_GOLD
        shape.line.fill.background()
        tf = shape.text_frame
        p = tf.add_paragraph()
        p.text = num
        p.alignment = PP_ALIGN.CENTER
        style_text_run(p.runs[0], 18, RGBColor(0,0,0), True)

        # Title
        txBox = slide.shapes.add_textbox(x, y, Inches(card_w), Inches(0.5))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        style_text_run(p.runs[0], 16, TEXT_PRIMARY, True)

        # Desc
        txBox = slide.shapes.add_textbox(x, y + Inches(0.5), Inches(card_w), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = desc
        p.alignment = PP_ALIGN.CENTER
        style_text_run(p.runs[0], 12, TEXT_SECONDARY)


    # --- Slide 5: Tech Stack ---
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    # Heading
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Technology Stack"
    p.alignment = PP_ALIGN.CENTER
    style_text_run(p.runs[0], 36, ACCENT_GOLD, True)

    techs = ["Flutter", "Python", "PyTorch", "YOLO", "Firebase", "Flask", "React"]
    
    # List Layout
    y_start = 2.5
    for i, tech in enumerate(techs):
        x = Inches(4)
        y = Inches(y_start + i * 0.5)
        
        # Bullet
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.1), Inches(0.15), Inches(0.15))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_GOLD
        shape.line.fill.background()
        
        # Text
        txBox = slide.shapes.add_textbox(x + Inches(0.3), y, Inches(4), Inches(0.5))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = tech
        style_text_run(p.runs[0], 18, TEXT_PRIMARY)

    prs.save('BatsmanPro_Presentation_Dark.pptx')
    print("Presentation saved successfully as BatsmanPro_Presentation_Dark.pptx!")

if __name__ == "__main__":
    create_presentation()
